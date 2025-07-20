import streamlit as st
import PyPDF2
import docx
import openpyxl
from pptx import Presentation
from huggingface_hub import InferenceClient
import time
import io

# Initialize DeepSeek client
client = InferenceClient()
model_id = "deepseek-ai/DeepSeek-V3-0324"

# ------------------------ #
#        Utilities         #
# ------------------------ #

def typewriter_effect(text, delay=0.01):
    output = ""
    for word in text.split():
        output += f"{word} "
        st.text(output)
        time.sleep(delay)

def extract_text_from_pdf(file):
    reader = PyPDF2.PdfReader(file)
    return "\n".join(page.extract_text() for page in reader.pages if page.extract_text())

def extract_text_from_docx(file):
    doc = docx.Document(file)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_from_xlsx(file):
    wb = openpyxl.load_workbook(file)
    text = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = [str(cell) for cell in row if cell is not None]
            text.append(" ".join(row_text))
    return "\n".join(text)

def extract_text(file, file_type):
    if file_type == "pdf":
        return extract_text_from_pdf(file)
    elif file_type == "docx":
        return extract_text_from_docx(file)
    elif file_type == "pptx":
        return extract_text_from_pptx(file)
    elif file_type == "xlsx":
        return extract_text_from_xlsx(file)
    else:
        return ""

def ask_deepseek(context, history, question):
    system_msg = (
        "You are a helpful assistant. Use the provided document context to answer user queries."
        " If the answer isn't found, say 'Not found in document'."
    )

    messages = [{"role": "system", "content": system_msg}]
    for q, a in history:
        messages.append({"role": "user", "content": q})
        messages.append({"role": "assistant", "content": a})

    messages.append({"role": "user", "content": f"Context:\n{context}\n\nQuestion: {question}"})

    response = client.chat.completions.create(
        model=model_id,
        messages=messages
    )
    return response.choices[0].message.content

def summarize_with_deepseek(text):
    prompt = f"Summarize the following document:\n\n{text}"
    messages = [{"role": "user", "content": prompt}]
    response = client.chat.completions.create(model=model_id, messages=messages)
    return response.choices[0].message.content


# ------------------------ #
#        QnA Page          #
# ------------------------ #

def qna_page():
    st.title("🧠 QnA with Uploaded Document")
    uploaded_files = st.file_uploader("Upload files", type=["pdf", "docx", "pptx", "xlsx"], accept_multiple_files=True)

    if "contexts" not in st.session_state:
        st.session_state.contexts = {}
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = {}
    if "active_file" not in st.session_state:
        st.session_state.active_file = None

    # Extract and store context for each uploaded file
    for file in uploaded_files:
        name = file.name
        file_ext = name.split(".")[-1].lower()
        if name not in st.session_state.contexts:
            text = extract_text(file, file_ext)
            st.session_state.contexts[name] = text
            st.session_state.chat_history[name] = []

    # File selector
    if st.session_state.contexts:
        st.session_state.active_file = st.selectbox("Choose a file to chat with", list(st.session_state.contexts.keys()))

    # Clear chat button
    if st.session_state.active_file and st.button("🗑️ Clear Chat"):
        st.session_state.chat_history[st.session_state.active_file] = []

    # Display chat interface
    if st.session_state.active_file:
        context = st.session_state.contexts[st.session_state.active_file]
        history = st.session_state.chat_history[st.session_state.active_file]

        for q, a in history:
            with st.chat_message("user"):
                st.markdown(q)
            with st.chat_message("assistant"):
                st.markdown(a)

        if user_prompt := st.chat_input("Ask a question about the document"):
            with st.chat_message("user"):
                st.markdown(user_prompt)

            with st.spinner("DeepSeek is thinking..."):
                answer = ask_deepseek(context, history, user_prompt)

            st.session_state.chat_history[st.session_state.active_file].append((user_prompt, answer))

            with st.chat_message("assistant"):
                st.markdown(answer)


# ------------------------ #
#      Summarize Page      #
# ------------------------ #

def summarise_page():
    st.title("📄 Summarize Document")

    uploaded_file = st.file_uploader("Upload a file to summarize", type=["pdf", "docx", "pptx", "xlsx"], key="summary_uploader")

    if uploaded_file:
        file_ext = uploaded_file.name.split(".")[-1].lower()
        text = extract_text(uploaded_file, file_ext)

        if st.button("Summarize"):
            with st.spinner("Generating summary..."):
                summary = summarize_with_deepseek(text)
            st.markdown("### ✨ Summary")
            st.write(summary)

# ------------------------ #
#         Main App         #
# ------------------------ #

st.sidebar.title("🧭 Navigation")
page = st.sidebar.radio("Go to", ["QnA", "Summarise"])

if page == "QnA":
    qna_page()
else:
    summarise_page()
